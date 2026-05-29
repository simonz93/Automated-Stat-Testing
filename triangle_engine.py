"""
sync_triangles.py  (v10)
========================
Applies ▲/▼ significance triangles from a survey data Excel file into a
PowerPoint deck.

EXCEL LAYOUT EXPECTED
---------------------
  Row 1-4 : Hierarchical column headers (Total/Country/Market → Region → Sub-col)
  Row 5   : "Wave" label row
  Row 6   : "Wave 4" / "Wave 5" labels
  Row 7   : Letter codes (A, B, C, D…)
  Row 8+  : Data with sig row immediately below each value row
  Column 1: Section headers (e.g. "D6_1. Agree Summary: …") and row labels.

The script finds the Wave 4 / Wave 5 column pair whose header ancestry
contains BOTH the configured region and sub-column.

SHAPE NAMING CONVENTIONS (set in PowerPoint Selection Pane)
-----------------------------------------------------------
Every shape to process starts with "x ". The script auto-scans for them.

1. CHART (single section)
     "x D6."           →  code=D6
     "x B2. (2)"       →  code=B2  (the "(2)" suffix is a variant tag, stripped)
     "x C1_1. Some descriptive label"  →  code=C1_1  (label after dot ignored)

2. CHART (multi-section, stacked / mixed sources)
     "x1 A2__1. x2 A2__2. x3 A2__3."
       → declares this chart pulls from sections A2__1, A2__2, A2__3.
         Each series in the chart is matched to a section by the section
         code that appears at the start of the series name in the chart data.

3. TEXTBOX — exact-label match
     "x B1. Net: Very/somewhat satisfied"
       → look up row label "Net: Very/somewhat satisfied" in section B1

4. TEXTBOX — rank lookup (new in v10)
     "x D6_1. 1"  → rank 1 (highest %) value-row in section D6_1
     "x D6_1. 2"  → rank 2 (second highest)
       The script sorts the section's value rows by the Wave 5 column
       (descending) and picks the Nth row. Triangle = sig of that row.

5. TABLE
     "x A2__3."  → treated like a chart whose categories are the table rows.

SIG LOGIC
---------
  Sig row is immediately below each data row.
  Letter (Wave 5's column letter) in Wave 4's sig col → ▲
  Letter (Wave 4's column letter) in Wave 5's sig col → ▼
  abs(round(w4*100) - round(w5*100)) < MIN_DIFF_PCT → no triangle

CHART LABEL POLICY
------------------
The script only modifies data labels that ALREADY EXIST as individual
dLbl entries in the chart. It never creates a new label slot. To get a
triangle on a bar, that bar must already have a data label visible in
the chart in PowerPoint.
"""

import functools, re, sys, uuid
from datetime import datetime
from pathlib import Path
from lxml import etree

try:
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from pptx import Presentation
except ImportError:
    sys.exit("Run:  pip install openpyxl python-pptx lxml")

# ── CONFIGURE ────────────────────────────────────────────────────────────────
EXCEL_FILE  = "test_data.xlsx"
PPTX_FILE   = "test.pptx"
OUTPUT_FILE = "test_updated.pptx"
REPORT_FILE = "sync_report.xlsx"
EXCEL_SHEET = "P_2_1"

# Default column cut applied to ALL shapes unless overridden below.
# Either set DEFAULT_CUT to an ordered list of header values (full hierarchy
# path), or leave it empty and the legacy DEFAULT_REGION / DEFAULT_SUB_COLUMN
# pair will be used.
DEFAULT_CUT        = []   # e.g. ["Country", "Rest of the world", "", "Total"]
DEFAULT_REGION     = "Rest of the world"
DEFAULT_SUB_COLUMN = "Total"

# Minimum rounded percentage-point difference to show a triangle
MIN_DIFF_PCT = 4

# ── CUSTOM COLUMN TARGETING ───────────────────────────────────────────────────
# Some shapes don't pull from the default region/sub-column (e.g. "Rest of the
# world > Total"). List those exceptions here. Each entry says: on this slide,
# this shape should read its data from a specific (region, sub_column) pair in
# the Excel header hierarchy instead of the default.
#
# The script finds the Wave 4 / Wave 5 column pair whose header ancestry contains
# BOTH the region and sub_column strings (matched against header rows 1-6).
#
# Format:
#   (slide_num, "shape_name", "region", "sub_column")
#
# Example header hierarchy in the data:
#   "All Bundle (NET)"  →  "Good Value"  →  Wave 4 / Wave 5   (cols I / J)
#   "All Bundle (NET)"  →  "Poor Value"  →  Wave 4 / Wave 5   (cols K / L)
SHAPE_OVERRIDES = [
    # Slide 26 — D5a/D5b read from the All Bundle (NET) sub-columns, not Total
    ( 26, "x D5a.", "All Bundle (NET)", "Good Value" ),   # cols I / J
    ( 26, "x D5b.", "All Bundle (NET)", "Poor Value" ),   # cols K / L
]

# Slides to skip entirely (e.g. ones still being designed)
SKIP_SLIDES = [42]

# Excel layout — change only if the data template changes
HEADER_ROWS = [1, 2, 3, 4, 6]   # rows scanned to match region/sub_column
LETTER_ROW  = 7                  # row containing sig letter codes (A, B, C…)
LABEL_COL   = 1                  # column containing row labels
# ─────────────────────────────────────────────────────────────────────────────

TRIANGLE_UP   = "▲"
TRIANGLE_DOWN = "▼"
COLOR_UP    = "3A9262"   # green for ▲
COLOR_DOWN  = "FF0000"   # red for ▼
COLOR_WHITE = "FFFFFF"   # white, for triangles over 100% stacked bars

# A ▲ on a green bar is forced white only when the bar is as green/dark as
# #64D246 or darker. We compare perceived luminance (ITU-R BT.601):
#   L = 0.299*R + 0.587*G + 0.114*B
# Reference #64D246 = RGB(100, 210, 70) → L ≈ 161.15. Any green bar with
# luminance <= this threshold is dark enough that a green triangle on it
# would be hard to read, so we switch to white.
_GREEN_REF_HEX = "64D246"
def _luminance(r, g, b):
    return 0.299 * r + 0.587 * g + 0.114 * b
_GREEN_BRIGHTNESS_MAX = _luminance(0x64, 0xD2, 0x46)   # ≈ 161.15
SHAPE_PREFIX  = "x"

C = "http://schemas.openxmlformats.org/drawingml/2006/chart"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

# In app/library mode these are not used; the engine works on in-memory bytes.
# When run as a CLI script they default to the configured filenames.
SCRIPT_DIR = Path(__file__).resolve().parent


# ── Small helpers ─────────────────────────────────────────────────────────────

@functools.lru_cache(maxsize=8192)
def normalise(text):
    """Lowercase, collapse whitespace, and fold smart quotes to ASCII."""
    t = str(text)
    t = t.replace("\u2018", "'").replace("\u2019", "'")
    t = t.replace("\u201c", '"').replace("\u201d", '"')
    t = t.replace("\u2013", "-").replace("\u2014", "-")
    return re.sub(r"[\s\xa0]+", " ", t).strip().lower()


def clean_sig(val):
    return str(val or "").strip().replace("*", "").upper()


def _pct_suffix(pct):
    """Return ' (W4: X% → W5: Y%)' from a (p4, p5) tuple, or '' if unavailable."""
    if not pct:
        return ""
    p4, p5 = pct
    parts = []
    if p4 is not None:
        parts.append(f"W4: {p4}%")
    if p5 is not None:
        parts.append(f"W5: {p5}%")
    return f" ({' → '.join(parts)})" if parts else ""


def pct_round(val):
    return round(float(val) * 100) if isinstance(val, (int, float)) else None


# ── Shape name parsing ────────────────────────────────────────────────────────

# Matches a single "x <CODE>." segment. The code can contain
# letters, digits, and underscores. Optional "(<digit>)" variant suffix.
SINGLE_CODE_RE = re.compile(
    r"^\s*x\s+([A-Za-z][\w]*)\.\s*(?:\((\d+)\)\s*\.?\s*)?(.*)$",
    re.IGNORECASE,
)
# xMULTI form — multi-source charts where each bar carries its section code
# as a prefix in the chart's category text (e.g. "D6_1. I get good...")
MULTI_PREFIX_RE = re.compile(
    r"^\s*xMULTI\s+([A-Za-z][\w]*)\.\s*$",
    re.IGNORECASE,
)
# Pattern at the start of a category label that exposes its section code
CAT_CODE_PREFIX_RE = re.compile(r"^([A-Za-z][\w]*)\.\s*", re.IGNORECASE)
RANK_LABEL_RE = re.compile(r"^(\d+)$")


def parse_shape_name(name):
    """
    Parse a shape name into a structured spec.
    Returns dict with:
      kind:       'single' | 'multi_root' | None
      codes:      list of declared section codes (always upper-case)
      root_code:  for 'multi_root' only — the root code (e.g. 'D6')
      label:      optional label string (None if no label or if it's a rank)
      rank:       int if label is a pure number, else None
    """
    if not name.startswith(SHAPE_PREFIX):
        return None

    # xMULTI form — single root code; per-bar sections come from category prefixes
    mm = MULTI_PREFIX_RE.match(name)
    if mm:
        return {
            "kind":      "multi_root",
            "codes":     [],          # resolved at runtime from chart categories
            "root_code": mm.group(1).upper(),
            "label":     None,
            "rank":      None,
        }

    # Single-section parse
    m = SINGLE_CODE_RE.match(name)
    if not m:
        return None
    code  = m.group(1).upper()
    label = (m.group(3) or "").strip() or None

    rank = None
    if label is not None:
        rm = RANK_LABEL_RE.match(label)
        if rm:
            rank  = int(rm.group(1))
            label = None  # rank replaces label
    return {
        "kind":  "single",
        "codes": [code],
        "label": label,
        "rank":  rank,
    }


# ── Excel column detection ────────────────────────────────────────────────────

def get_col_ancestry(ws, col, _cache={}):
    """Walk leftward across each HEADER_ROW to find the column's hierarchy."""
    if col in _cache:
        return _cache[col]
    labels = set()
    for hrow in HEADER_ROWS:
        for c in range(col, 0, -1):
            v = ws.cell(row=hrow, column=c).value
            if v is not None:
                labels.add(str(v).strip())
                break
    _cache[col] = labels
    return labels


def get_col_path(ws, col):
    """
    Ordered list of header values above a column, one per HEADER_ROW that sits
    ABOVE the wave row (filled leftward). e.g. ['Country', 'Rest of the world',
    '', 'Total']. Blank levels are kept as '' so positions stay aligned.
    """
    wave_row = find_wave_row(ws)
    levels = [r for r in HEADER_ROWS if wave_row is None or r < wave_row]
    path = []
    for hr in levels:
        val = ""
        for c in range(col, 0, -1):
            cell = ws.cell(row=hr, column=c).value
            if cell is not None:
                val = str(cell).strip()
                break
        path.append(val)
    return path


def build_header_tree(ws):
    """
    Build the hierarchy of column headers above the wave row, for cascading UI
    dropdowns. Returns:
      {
        "levels": N,                      # number of header rows above wave row
        "paths":  [ (path_tuple, col), ]  # every wave-column's full path
      }
    A 'path' is a tuple of the level values, e.g.
      ('Country', 'Rest of the world', '', 'Total').
    Only one entry per distinct path is kept (Wave 4 column is enough to map it).
    """
    wave_row = find_wave_row(ws)
    levels = [r for r in HEADER_ROWS if wave_row is None or r < wave_row]
    seen = set()
    paths = []
    for col in range(2, ws.max_column + 1):
        v = ws.cell(row=wave_row, column=col).value if wave_row else None
        if not v or str(v).strip() not in ("Wave 4", "Wave 5"):
            continue
        p = tuple(get_col_path(ws, col))
        if p in seen:
            continue
        seen.add(p)
        paths.append((p, col))
    return {"levels": len(levels), "paths": paths}


def find_wave_columns(ws, *cuts):
    """
    Locate the Wave 4 / Wave 5 column pair matching the given header values.

    Accepts either:
      • find_wave_columns(ws, region, sub_col)         (legacy 2-value form)
      • find_wave_columns(ws, lvl1, lvl2, lvl3, lvl4)   (full hierarchy path)
      • find_wave_columns(ws, ["lvl1", "lvl2", ...])    (path as a single list)

    Blank path entries ('') are ignored when matching. All non-blank values must
    appear in the column's header ancestry.
    """
    if len(cuts) == 1 and isinstance(cuts[0], (list, tuple)):
        cuts = tuple(cuts[0])
    match_set = {str(c).strip() for c in cuts if str(c).strip()}

    wave_row = find_wave_row(ws)
    if wave_row is None:
        sys.exit("ERROR: Could not find Wave 4 / Wave 5 label row in HEADER_ROWS")

    w4 = w5 = None
    for col in range(2, ws.max_column + 1):
        v = ws.cell(row=wave_row, column=col).value
        if v is None:
            continue
        wave_val = str(v).strip()
        if wave_val not in ("Wave 4", "Wave 5"):
            continue
        if match_set.issubset(get_col_ancestry(ws, col)):
            if wave_val == "Wave 4" and w4 is None:
                w4 = col
            elif wave_val == "Wave 5" and w5 is None:
                w5 = col
        if w4 and w5:
            break

    if not w4 or not w5:
        sys.exit(
            f"ERROR: No Wave 4 / Wave 5 column pair matching {sorted(match_set)}"
        )
    lw4 = clean_sig(ws.cell(row=LETTER_ROW, column=w4).value)
    lw5 = clean_sig(ws.cell(row=LETTER_ROW, column=w5).value)
    return w4, w5, lw4, lw5



# ── Header value scanning (for UI dropdowns) ──────────────────────────────────

def find_wave_row(ws):
    """Return the row number within HEADER_ROWS that holds the Wave 4/5 labels."""
    for r in HEADER_ROWS:
        for c in range(1, ws.max_column + 1):
            v = ws.cell(row=r, column=c).value
            if v and str(v).strip() in ("Wave 4", "Wave 5"):
                return r
    return None


def scan_header_values(ws):
    """
    Inspect the column-header hierarchy above each Wave 4/Wave 5 column and
    return the distinct values that appear at each level.

    Returns a dict:
      {
        "regions":   [sorted list of level-1/level-2 group labels],
        "sub_cols":  [sorted list of the labels sitting directly above the
                      wave row, i.e. the sub-column level],
        "sheets":    handled by the caller,
      }
    'regions' and 'sub_cols' are deliberately generous supersets so the UI can
    offer every label that could legitimately be typed into find_wave_columns.
    """
    wave_row = find_wave_row(ws)
    if wave_row is None:
        return {"regions": [], "sub_cols": []}

    # The sub-column level is the header row immediately above the wave row.
    # Everything above that contributes "region"/group candidates.
    rows_above = [r for r in HEADER_ROWS if r < wave_row]
    sub_row = max(rows_above) if rows_above else None
    upper_rows = [r for r in rows_above if r != sub_row]

    sub_vals = set()
    region_vals = set()

    # Walk the wave columns and collect the ancestry value at each level.
    for col in range(2, ws.max_column + 1):
        v = ws.cell(row=wave_row, column=col).value
        if not v or str(v).strip() not in ("Wave 4", "Wave 5"):
            continue
        # Sub-column: nearest non-empty cell at sub_row scanning leftward
        if sub_row is not None:
            for c in range(col, 0, -1):
                cell = ws.cell(row=sub_row, column=c).value
                if cell is not None:
                    sub_vals.add(str(cell).strip())
                    break
        # Region/group levels: nearest non-empty at each upper row
        for hr in upper_rows:
            for c in range(col, 0, -1):
                cell = ws.cell(row=hr, column=c).value
                if cell is not None:
                    region_vals.add(str(cell).strip())
                    break

    drop_blank = lambda s: [x for x in sorted(s) if x]
    return {"regions": drop_blank(region_vals), "sub_cols": drop_blank(sub_vals)}


# ── Excel section scanning ────────────────────────────────────────────────────

SECTION_CODE_RE = re.compile(r"^([A-Za-z][\w]*)\.\s*(.*)", re.IGNORECASE)


def find_all_sections(ws):
    """Return list of (code_upper, start_row, header_text) for every section."""
    sections = []
    for r in range(1, ws.max_row + 1):
        val = ws.cell(row=r, column=LABEL_COL).value
        if val and isinstance(val, str):
            m = SECTION_CODE_RE.match(val.strip())
            if m:
                sections.append((m.group(1).upper(), r, val.strip()))
    return sections


def section_rows(ws, start, end, w4, w5, lw4, lw5):
    """
    Walk rows (start+1..end) and return list of:
      (row_num, label, w4_val, w5_val, p4, p5, triangle_or_None)
    Skips Base/Unweighted base rows automatically.
    """
    end = end or ws.max_row + 1
    out = []
    for r in range(start + 1, end):
        label = ws.cell(row=r, column=LABEL_COL).value
        if label is None or not isinstance(label, str):
            continue
        if normalise(label) in ("base", "unweighted base") or normalise(label).startswith("base:"):
            continue
        w4v = ws.cell(row=r, column=w4).value
        w5v = ws.cell(row=r, column=w5).value
        p4, p5 = pct_round(w4v), pct_round(w5v)
        sig_w5 = clean_sig(ws.cell(row=r + 1, column=w5).value)
        sig_w4 = clean_sig(ws.cell(row=r + 1, column=w4).value)
        tri = None
        if p4 is not None and p5 is not None and abs(p5 - p4) >= MIN_DIFF_PCT:
            if sig_w5 == lw4:
                tri = TRIANGLE_UP
            elif sig_w4 == lw5:
                tri = TRIANGLE_DOWN
        out.append((r, label, w4v, w5v, p4, p5, tri))
    return out


def build_section_data(ws, sections, w4, w5, lw4, lw5):
    """
    Return dict: code → {
        'rows'        : list of (row_num, label, w4v, w5v, p4, p5, triangle),
        'by_label'    : { normalised_label : triangle },
        'ranked_w5'   : rows sorted by p5 descending (only rows with p5)
    }
    First occurrence of a code wins; later duplicates merge unique labels in.
    """
    result = {}
    for i, (code, start, _) in enumerate(sections):
        end = sections[i + 1][1] if i + 1 < len(sections) else None
        rows = section_rows(ws, start, end, w4, w5, lw4, lw5)
        by_label = {normalise(lbl): tri for _, lbl, _, _, _, _, tri in rows if tri}
        pct_by_label = {normalise(lbl): (p4, p5) for _, lbl, _, _, p4, p5, _ in rows}
        ranked_w5 = sorted(
            [(r_, lbl, w4v, w5v, p4, p5, tri) for r_, lbl, w4v, w5v, p4, p5, tri in rows if p5 is not None],
            key=lambda x: x[5],
            reverse=True,
        )
        if code not in result:
            result[code] = {"rows": rows, "by_label": by_label,
                            "pct_by_label": pct_by_label, "ranked_w5": ranked_w5}
        else:
            # Merge unique sig labels (preserves "first wins" while not losing duplicates)
            for k, v in by_label.items():
                if k not in result[code]["by_label"]:
                    result[code]["by_label"][k] = v
            for k, v in pct_by_label.items():
                if k not in result[code]["pct_by_label"]:
                    result[code]["pct_by_label"][k] = v
    return result


def _default_cut():
    """The default column cut as a tuple of header values."""
    cut = [c for c in globals().get("DEFAULT_CUT", []) if c]
    if cut:
        return tuple(cut)
    # Legacy fallback: region + sub-column
    return tuple(v for v in (DEFAULT_REGION, DEFAULT_SUB_COLUMN) if v)


def _override_cut(entry):
    """
    Normalise a SHAPE_OVERRIDES entry to (slide, shape_name, cut_tuple).
    Accepts:
      (slide, name, region, sub)        legacy 4-tuple
      (slide, name, [lvl1, lvl2, ...])  path form (3-tuple with a list)
    """
    if len(entry) == 4:
        slide, name, region, sub = entry
        cut = tuple(v for v in (str(region).strip(), str(sub).strip()) if v)
    else:
        slide, name, path = entry
        cut = tuple(str(v).strip() for v in path if str(v).strip())
    return int(slide), name, cut


def load_excel_data(ws, log=print):
    """Build section data for every distinct column cut in use (path tuples)."""
    all_sections = find_all_sections(ws)

    default_cut = _default_cut()
    col_specs = {default_cut}
    for entry in SHAPE_OVERRIDES:
        _, _, cut = _override_cut(entry)
        col_specs.add(cut)

    data_by_col = {}
    for cut in col_specs:
        w4, w5, lw4, lw5 = find_wave_columns(ws, list(cut))
        tag = "" if cut == default_cut else " [override]"
        log(f"  cols -> W4=col{w4}('{lw4}') W5=col{w5}('{lw5}') | "
            f"{' > '.join(cut)}{tag}")
        data_by_col[cut] = build_section_data(ws, all_sections, w4, w5, lw4, lw5)

    log(f"  {len(all_sections)} sections discovered.")
    return data_by_col, all_sections


# ── Chart XML helpers ─────────────────────────────────────────────────────────

def chart_categories(series):
    """Return [(idx, text), …] of categories for a series."""
    try:
        cat_el = series._element.find(f".//{{{C}}}cat")
        if cat_el is None:
            return []
        cache = cat_el.find(f".//{{{C}}}strRef/{{{C}}}strCache")
        if cache is None:
            return []
        result = []
        for p in cache.findall(f"{{{C}}}pt"):
            v = p.find(f"{{{C}}}v")
            if v is not None and v.text is not None:
                result.append((int(p.get("idx", 0)), v.text))
        return result
    except Exception:
        return []


def find_existing_dlbl(dLbls, idx):
    """
    Return a dLbl element for idx, ready to receive a triangle.
    Behaviour:
      1. If an individual <c:dLbl> for idx already exists → return it.
      2. Else if the series-level <c:dLbls> has <c:showVal val="1"/>, all bars
         in this series already show labels by default — create an individual
         dLbl entry as an override (PowerPoint uses it in place of the default).
      3. Else (no series-level showVal=1 and no individual entry) → return None,
         because creating a new label would add one where there was none.
    """
    if dLbls is None:
        return None

    # 1. Existing individual entry
    for dLbl in dLbls.findall(f"{{{C}}}dLbl"):
        idx_el = dLbl.find(f"{{{C}}}idx")
        if idx_el is not None and int(idx_el.get("val")) == idx:
            return dLbl

    # 2. Series shows labels by default? Check showVal=1 at series level
    showVal = dLbls.find(f"{{{C}}}showVal")
    series_shows_all = showVal is not None and showVal.get("val") == "1"
    if not series_shows_all:
        return None

    # Create a new dLbl that inherits the series default by including a
    # placeholder [VALUE] field, then we'll append the triangle.
    # Insert in sorted-by-idx position so PowerPoint reads it in order.
    new_dLbl = etree.Element(f"{{{C}}}dLbl")
    etree.SubElement(new_dLbl, f"{{{C}}}idx").set("val", str(idx))
    # Build the tx with [VALUE] field — preserves the series-level formatting
    tx   = etree.SubElement(new_dLbl, f"{{{C}}}tx")
    rich = etree.SubElement(tx,    f"{{{C}}}rich")
    bp   = etree.SubElement(rich,  f"{{{A}}}bodyPr")
    bp.set("wrap", "square")
    etree.SubElement(rich, f"{{{A}}}lstStyle")
    p    = etree.SubElement(rich,  f"{{{A}}}p")
    fld  = etree.SubElement(p,     f"{{{A}}}fld",
                             attrib={"id": "{" + str(uuid.uuid4()).upper() + "}",
                                     "type": "VALUE"})
    etree.SubElement(fld, f"{{{A}}}rPr",
                     attrib={"lang": "en-US", "smtClean": "0"})
    etree.SubElement(fld, f"{{{A}}}pPr")
    etree.SubElement(fld, f"{{{A}}}t").text = "[VALUE]"
    etree.SubElement(p, f"{{{A}}}endParaRPr").set("lang", "en-US")
    # Series-level display flags inherited; explicitly mirror showVal so the
    # override doesn't accidentally hide the value when we add the triangle.
    for tag, val in [("showLegendKey", "0"), ("showVal", "1"),
                     ("showCatName", "0"), ("showSerName", "0"),
                     ("showPercent", "0"), ("showBubbleSize", "0")]:
        etree.SubElement(new_dLbl, f"{{{C}}}{tag}").set("val", val)

    # Insert into dLbls in idx order, before any series-level config elements
    inserted = False
    for child in list(dLbls):
        c_idx = child.find(f"{{{C}}}idx") if child.tag.endswith("}dLbl") else None
        if c_idx is not None and int(c_idx.get("val")) > idx:
            child.addprevious(new_dLbl)
            inserted = True
            break
    if not inserted:
        # Insert before the first non-dLbl child (e.g. spPr / txPr / showLegendKey)
        first_non_dlbl = None
        for child in list(dLbls):
            if not child.tag.endswith("}dLbl"):
                first_non_dlbl = child
                break
        if first_non_dlbl is not None:
            first_non_dlbl.addprevious(new_dLbl)
        else:
            dLbls.append(new_dLbl)
    return new_dLbl


# Cache of theme scheme colours: {scheme_name_lower: "RRGGBB"}
_THEME_COLORS = {}

def load_theme_colors(prs):
    """Populate _THEME_COLORS from the first slide master's theme."""
    global _THEME_COLORS
    _THEME_COLORS = {}
    try:
        master = prs.slide_masters[0]
        theme_part = None
        for rel in master.part.rels.values():
            if "theme" in rel.reltype:
                theme_part = rel.target_part
                break
        if theme_part is None:
            return
        root = etree.fromstring(theme_part.blob)
        clrScheme = root.find(f".//{{{A}}}clrScheme")
        if clrScheme is None:
            return
        for child in clrScheme:
            tag = child.tag.split("}")[-1].lower()
            srgb = child.find(f"{{{A}}}srgbClr")
            sysc = child.find(f"{{{A}}}sysClr")
            if srgb is not None:
                _THEME_COLORS[tag] = srgb.get("val")
            elif sysc is not None and sysc.get("lastClr"):
                _THEME_COLORS[tag] = sysc.get("lastClr")
        # OOXML quirk: dk1/lt1 map to tx1/bg1, dk2/lt2 to tx2/bg2
        for a, b in [("dk1", "tx1"), ("lt1", "bg1"), ("dk2", "tx2"), ("lt2", "bg2")]:
            if a in _THEME_COLORS and b not in _THEME_COLORS:
                _THEME_COLORS[b] = _THEME_COLORS[a]
    except Exception:
        pass


def _apply_color_mods(r, g, b, clr_el):
    """
    Apply OOXML colour transform child elements (lumMod, lumOff, shade, tint)
    to an (r,g,b) triple. These lighten/darken theme colours — without them a
    pale tint of a dark theme green would be misread as the dark base green.
    Approximations in RGB space; close enough for green/red classification.
    """
    def pct(el, name):
        child = el.find(f"{{{A}}}{name}")
        if child is not None and child.get("val") is not None:
            try:
                return int(child.get("val")) / 100000.0
            except ValueError:
                return None
        return None

    lum_mod = pct(clr_el, "lumMod")
    lum_off = pct(clr_el, "lumOff")
    shade   = pct(clr_el, "shade")
    tint    = pct(clr_el, "tint")

    # lumMod scales luminance, lumOff adds luminance (both as fractions of 255)
    if lum_mod is not None:
        r, g, b = r * lum_mod, g * lum_mod, b * lum_mod
    if lum_off is not None:
        add = lum_off * 255
        r, g, b = r + add, g + add, b + add
    # shade darkens toward black, tint lightens toward white
    if shade is not None:
        r, g, b = r * shade, g * shade, b * shade
    if tint is not None:
        r = r * tint + 255 * (1 - tint)
        g = g * tint + 255 * (1 - tint)
        b = b * tint + 255 * (1 - tint)

    clamp = lambda v: max(0, min(255, int(round(v))))
    return clamp(r), clamp(g), clamp(b)


def resolve_fill_hex(spPr):
    """
    Resolve a shape/series fill to a 'RRGGBB' hex string, applying any
    luminance/shade/tint modifiers. Returns None if no solid fill found.
    """
    if spPr is None:
        return None
    srgb = spPr.find(f".//{{{A}}}srgbClr")
    scheme = spPr.find(f".//{{{A}}}schemeClr")

    if srgb is not None:
        base = srgb.get("val")
        clr_el = srgb
    elif scheme is not None:
        base = _THEME_COLORS.get(scheme.get("val", "").lower())
        clr_el = scheme
    else:
        return None

    if not base or len(base) != 6:
        return base
    try:
        r = int(base[0:2], 16)
        g = int(base[2:4], 16)
        b = int(base[4:6], 16)
    except ValueError:
        return base

    r, g, b = _apply_color_mods(r, g, b, clr_el)
    return f"{r:02X}{g:02X}{b:02X}"


def classify_color(hex_color):
    """Roughly classify a hex colour as 'green', 'red', or 'other'."""
    if not hex_color or len(hex_color) != 6:
        return "other"
    try:
        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
    except ValueError:
        return "other"
    # Green: green channel dominant AND as dark as #64D246 or darker
    if g > r + 25 and g > b + 25 and _luminance(r, g, b) <= _GREEN_BRIGHTNESS_MAX:
        return "green"
    # Red: red channel clearly dominant and not orange/yellow
    if r > g + 60 and r > b + 60 and g < 120:
        return "red"
    return "other"


# dLblPos values that place the label INSIDE the bar
_INSIDE_POSITIONS = {"ctr", "inend", "inbase"}

def get_effective_dlbl_pos(dLbl, series_dLbls):
    """Resolve a label's position: own dLblPos → series dLbls dLblPos → None."""
    for src in (dLbl, series_dLbls):
        if src is None:
            continue
        pos_el = src.find(f"{{{C}}}dLblPos")
        if pos_el is not None and pos_el.get("val"):
            return pos_el.get("val").lower()
    return None


def should_force_white(triangle, series_el, dLbl, series_dLbls):
    """
    True if the triangle would clash with the bar it sits on:
    label is inside the bar (ctr/inEnd/inBase) AND
      (▲ on a green bar) or (▼ on a red bar).
    """
    pos = get_effective_dlbl_pos(dLbl, series_dLbls)
    if pos not in _INSIDE_POSITIONS:
        return False
    spPr = series_el.find(f"{{{C}}}spPr")
    bar_class = classify_color(resolve_fill_hex(spPr))
    if triangle == TRIANGLE_UP and bar_class == "green":
        return True
    if triangle == TRIANGLE_DOWN and bar_class == "red":
        return True
    return False


def triangle_color(triangle, white=False):
    """Return the hex colour for a triangle. White overrides direction colour."""
    if white:
        return COLOR_WHITE
    return COLOR_UP if triangle == TRIANGLE_UP else COLOR_DOWN


def append_triangle_to_dlbl(dLbl, triangle, white=False):
    """
    Append a triangle run to an existing dLbl's text body, preserving the
    original structure. Returns True on success. Reorders <c:tx> to be
    the second child (right after <c:idx>) per OOXML spec.
    """
    tx = dLbl.find(f"{{{C}}}tx")
    if tx is None:
        return False
    # rich can be in either the c: or a: namespace
    rich = tx.find(f".//{{{C}}}rich")
    if rich is None:
        rich = tx.find(f".//{{{A}}}rich")
    if rich is None:
        return False
    paras = rich.findall(f"{{{A}}}p")
    if not paras:
        return False
    p = paras[-1]
    # Strip any old triangle runs we may have inserted previously
    for r in list(p.findall(f"{{{A}}}r")):
        t = r.find(f"{{{A}}}t")
        if t is not None and t.text and t.text.strip() in (TRIANGLE_UP, TRIANGLE_DOWN):
            p.remove(r)
    # Insert the new triangle run before endParaRPr if present
    end_rpr = p.find(f"{{{A}}}endParaRPr")
    r_tri = etree.Element(f"{{{A}}}r")
    rPr = etree.SubElement(r_tri, f"{{{A}}}rPr", attrib={"lang": "en-US", "smtClean": "0"})
    # Colour the triangle
    fill = etree.SubElement(rPr, f"{{{A}}}solidFill")
    etree.SubElement(fill, f"{{{A}}}srgbClr").set("val", triangle_color(triangle, white))
    etree.SubElement(r_tri, f"{{{A}}}t").text = " " + triangle
    if end_rpr is not None:
        end_rpr.addprevious(r_tri)
    else:
        p.append(r_tri)
    # Ensure wrap=none on bodyPr (matches the chart's own convention)
    bp = rich.find(f"{{{A}}}bodyPr")
    if bp is not None:
        bp.set("wrap", "none")
    # Ensure <c:tx> sits at position 1 (after <c:idx>); out-of-order tx
    # causes PowerPoint to double-render the label.
    children = list(dLbl)
    if len(children) >= 2 and children[1] is not tx:
        # Rebuild order: idx, tx, then everything else
        idx_el = next((c for c in children if c.tag.endswith("}idx")), None)
        for child in children:
            dLbl.remove(child)
        if idx_el is not None:
            dLbl.append(idx_el)
        dLbl.append(tx)
        for child in children:
            if child is not idx_el and child is not tx:
                dLbl.append(child)
    return True


# ── Text writers ──────────────────────────────────────────────────────────────

def _set_run_color(run, hex_color):
    from pptx.dml.color import RGBColor
    run.font.color.rgb = RGBColor.from_string(hex_color)


def write_triangle_to_textbox(shape, triangle, white=False):
    """
    Append a triangle as its OWN run (so it can be coloured independently of
    the surrounding text) to the end of the textbox.
    """
    tf = shape.text_frame
    last_para = tf.paragraphs[-1]
    # Strip any previous triangle that may sit at the tail of the last run
    if last_para.runs:
        tail = last_para.runs[-1]
        tail.text = tail.text.rstrip("▲▼ ").rstrip()
    # Add the triangle as a separate run with its own colour
    tri_run = last_para.add_run()
    tri_run.text = " " + triangle
    _set_run_color(tri_run, triangle_color(triangle, white))


def _copy_font_from_first_run(tf):
    """Capture font properties from the first run in a text frame, if any."""
    props = {}
    if tf.paragraphs and tf.paragraphs[0].runs:
        r0 = tf.paragraphs[0].runs[0]
        props = {"bold": r0.font.bold, "italic": r0.font.italic, "size": r0.font.size}
        try:
            props["color"] = r0.font.color.rgb if r0.font.color.type else None
        except Exception:
            props["color"] = None
    return props


def _apply_font(run, props):
    if props.get("bold")   is not None: run.font.bold   = props["bold"]
    if props.get("italic") is not None: run.font.italic = props["italic"]
    if props.get("size")   is not None: run.font.size   = props["size"]
    if props.get("color")  is not None:
        run.font.color.rgb = props["color"]


def set_table_cell_triangle(cell, triangle, white=False):
    """Write a triangle into an (empty) cell, inheriting nearby font style."""
    tf = cell.text_frame
    props = _copy_font_from_first_run(tf)
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""
    run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
    run.text = triangle
    _apply_font(run, props)
    _set_run_color(run, triangle_color(triangle, white))


def append_triangle_to_table_cell(cell, triangle, white=False):
    """Append a triangle (its own coloured run) after the cell's existing text."""
    tf = cell.text_frame
    last_para = tf.paragraphs[-1]
    if last_para.runs:
        tail = last_para.runs[-1]
        tail.text = tail.text.rstrip("▲▼ ").rstrip()
    tri_run = last_para.add_run()
    tri_run.text = " " + triangle
    _set_run_color(tri_run, triangle_color(triangle, white))


# ── Chart processing ──────────────────────────────────────────────────────────

def process_chart(shape, spec, section_data):
    """
    Apply triangles to a chart's existing data labels.
    Returns (updates_count, list_of_log_strings, list_of_item_strings).
    """
    chart = shape.chart
    logs  = []
    items = []
    n     = 0
    kind  = spec["kind"]

    # ── multi_root: xMULTI form ──────────────────────────────────────────────
    # Two layouts handled:
    #   A) Categories carry the section code (e.g. cats="D6_1. Statement…")
    #      → series name = row label, lookup in cat's section
    #   B) Series names carry the section code (e.g. series="A2__1. Aware Summary")
    #      → category text = row label, lookup in series' section
    if kind == "multi_root":
        root = spec["root_code"]
        root_u = root.upper()

        # Snapshot all series with their categories
        series_list = []
        for series in chart.series:
            dLbls = series._element.find(f"{{{C}}}dLbls")
            if dLbls is None:
                continue
            cats = chart_categories(series)
            series_list.append((series, dLbls, cats, series.name or ""))

        # Decide layout: do series names carry the root prefix?
        series_carry_code = any(
            normalise(sn).startswith(normalise(root_u + "_")) or
            normalise(sn).startswith(normalise(root_u + "."))
            for _, _, _, sn in series_list
        )

        if series_carry_code:
            # Layout B: series → section, category → row label
            for series, dLbls, cats, series_name in series_list:
                cm = CAT_CODE_PREFIX_RE.match(series_name)
                if not cm:
                    continue
                ser_code = cm.group(1).upper()
                if not ser_code.startswith(root_u):
                    continue
                sec = section_data.get(ser_code, {})
                if not sec:
                    continue
                label_map = sec.get("by_label", {})
                pct_map   = sec.get("pct_by_label", {})
                for cat_idx, cat_text in cats:
                    tri = label_map.get(normalise(cat_text))
                    if not tri:
                        continue
                    dLbl = find_existing_dlbl(dLbls, cat_idx)
                    if dLbl is None:
                        continue
                    white = should_force_white(tri, series._element, dLbl, dLbls)
                    if append_triangle_to_dlbl(dLbl, tri, white):
                        logs.append(
                            f"chart '{shape.name}' [{ser_code}] cat[{cat_idx}] '{cat_text[:40]}' → {tri}"
                        )
                        items.append(f"[{ser_code}] '{cat_text[:40]}' → {tri}"
                                     f"{_pct_suffix(pct_map.get(normalise(cat_text)))}")
                        n += 1
            return n, logs, items

        # Layout A: category → section, series name → row label
        for series, dLbls, cats, series_name in series_list:
            for cat_idx, cat_text in cats:
                cm = CAT_CODE_PREFIX_RE.match(cat_text)
                if not cm:
                    continue
                cat_code = cm.group(1).upper()
                if not cat_code.startswith(root_u):
                    continue
                sec = section_data.get(cat_code, {})
                if not sec:
                    continue
                label_map = sec.get("by_label", {})
                pct_map   = sec.get("pct_by_label", {})
                tri = label_map.get(normalise(series_name))
                if not tri:
                    continue
                dLbl = find_existing_dlbl(dLbls, cat_idx)
                if dLbl is None:
                    continue
                white = should_force_white(tri, series._element, dLbl, dLbls)
                if append_triangle_to_dlbl(dLbl, tri, white):
                    logs.append(
                        f"chart '{shape.name}' [{cat_code}] series='{series_name[:25]}' "
                        f"cat[{cat_idx}] → {tri}"
                    )
                    items.append(f"[{cat_code}] '{series_name[:25]}' → {tri}"
                                 f"{_pct_suffix(pct_map.get(normalise(series_name)))}")
                    n += 1
        return n, logs, items

    # Resolve which sections this chart pulls from
    codes = spec["codes"]

    for series in chart.series:
        series_el = series._element
        dLbls = series_el.find(f"{{{C}}}dLbls")
        if dLbls is None:
            continue
        cats = chart_categories(series)
        series_name = series.name or ""

        # Single-code chart
        code = codes[0]
        sec  = section_data.get(code, {})
        label_map = sec.get("by_label", {})
        pct_map   = sec.get("pct_by_label", {})

        # Mode A: the series name matches a row label (e.g. "Very satisfied")
        sname_tri = label_map.get(normalise(series_name))
        if sname_tri and cats:
            # Series name = the metric; apply to every category that has an
            # existing dLbl (we never create new ones).
            series_n = 0
            for cat_idx, cat_text in cats:
                dLbl = find_existing_dlbl(dLbls, cat_idx)
                if dLbl is None:
                    continue
                white = should_force_white(sname_tri, series._element, dLbl, dLbls)
                if append_triangle_to_dlbl(dLbl, sname_tri, white):
                    logs.append(f"chart '{shape.name}' series='{series_name}' cat[{cat_idx}] → {sname_tri}")
                    series_n += 1
            if series_n:
                items.append(f"'{series_name}' → {sname_tri}"
                             f"{_pct_suffix(pct_map.get(normalise(series_name)))}")
                n += series_n
            continue

        # Mode B: categories are the row labels (single-series bar chart)
        if cats:
            for cat_idx, cat_text in cats:
                # Try the raw category text first
                tri = label_map.get(normalise(cat_text))
                # If no match, try stripping a leading "<code>. " section prefix
                # (e.g. "C7_13. Battlegrounds" → "Battlegrounds")
                if not tri:
                    stripped = CAT_CODE_PREFIX_RE.sub("", cat_text, count=1)
                    if stripped != cat_text:
                        tri = label_map.get(normalise(stripped))
                if not tri:
                    continue
                dLbl = find_existing_dlbl(dLbls, cat_idx)
                if dLbl is None:
                    continue
                white = should_force_white(tri, series._element, dLbl, dLbls)
                if append_triangle_to_dlbl(dLbl, tri, white):
                    logs.append(f"chart '{shape.name}' cat[{cat_idx}] '{cat_text[:45]}' → {tri}")
                    _norm = normalise(cat_text)
                    _pct  = pct_map.get(_norm) or pct_map.get(
                        normalise(CAT_CODE_PREFIX_RE.sub("", cat_text, count=1)))
                    items.append(f"'{cat_text[:45]}' → {tri}{_pct_suffix(_pct)}")
                    n += 1

    return n, logs, items


# ── Table processing ──────────────────────────────────────────────────────────

def process_table(shape, spec, section_data):
    """
    Apply triangles to a table.
    Column 0 holds row labels, the value column holds percentages.
    If the table has a dedicated empty trailing column, the triangle is
    written there (preserving the % value). Otherwise the triangle is
    appended to the value cell text rather than replacing it.
    """
    code = spec["codes"][0]
    sec  = section_data.get(code, {})
    label_map = sec.get("by_label", {})
    pct_map   = sec.get("pct_by_label", {})
    table = shape.table
    ncols = len(table.columns)
    nrows = len(table.rows)
    logs  = []
    items = []
    n     = 0

    # The triangle is appended directly after the "%" in the value cell.
    # Find the value column: the first column (after col 0) whose cells
    # contain a "%". Defaults to column 1.
    value_col = 1
    for c in range(1, ncols):
        if any("%" in table.cell(r, c).text for r in range(nrows)):
            value_col = c
            break

    for r in range(nrows):
        row_label = table.cell(r, 0).text.strip()
        if not row_label:
            continue
        tri = label_map.get(normalise(row_label))
        if not tri:
            # Try stripping a section-code prefix
            stripped = CAT_CODE_PREFIX_RE.sub("", row_label, count=1)
            if stripped != row_label:
                tri = label_map.get(normalise(stripped))
        if not tri:
            continue
        # Append the triangle right after the % in the value cell
        append_triangle_to_table_cell(table.cell(r, value_col), tri)
        logs.append(f"table '{shape.name}' row {r} '{row_label[:40]}' → {tri}")
        _pct = pct_map.get(normalise(row_label)) or pct_map.get(
            normalise(CAT_CODE_PREFIX_RE.sub("", row_label, count=1)))
        items.append(f"'{row_label[:40]}' → {tri}{_pct_suffix(_pct)}")
        n += 1
    return n, logs, items


# ── Shape dispatcher ──────────────────────────────────────────────────────────

STATUS_OK      = "Data found – stat testing applied"
STATUS_NO_SIG  = "Data found – no stat testing"
STATUS_WARN    = "Data not found"
STATUS_MISSING = "Data not found"


def chart_grouping(shape):
    """Return the chart's grouping value (e.g. 'percentStacked') or None."""
    if not shape.has_chart:
        return None
    try:
        for el in shape.chart._chartSpace.iter():
            if el.tag.split("}")[-1] == "grouping":
                return el.get("val")
    except Exception:
        pass
    return None


def _shapes_overlap(a, b):
    try:
        return (a.left < b.left + b.width and a.left + a.width > b.left and
                a.top  < b.top  + b.height and a.top  + a.height > b.top)
    except Exception:
        return False


def textbox_over_pct_stacked(slide, textbox):
    """True if the textbox overlaps a 100% (percentStacked) bar chart on the slide."""
    for other in slide.shapes:
        if other is textbox or not other.has_chart:
            continue
        if chart_grouping(other) == "percentStacked" and _shapes_overlap(textbox, other):
            return True
    return False


def process_shape(slide_num, slide, shape, spec, section_data, cut, report_rows):
    """Dispatch one shape to chart/table/textbox handling. Returns updates count."""
    col_label = " > ".join(cut) if cut else "(default)"
    base = {"slide": slide_num, "shape_name": shape.name}

    if spec is None:
        report_rows.append({**base, "status": STATUS_WARN, "items_applied": "",
                            "triangles_added": 0})
        print(f"  [SKIP] '{shape.name}': unparseable name")
        return 0

    # Validate declared sections exist (multi_root resolves codes at runtime
    # from category prefixes, so skip this check for that kind)
    if spec.get("kind") != "multi_root":
        missing = [c for c in spec["codes"] if c not in section_data]
        if missing:
            report_rows.append({**base, "status": STATUS_MISSING, "items_applied": "",
                                "triangles_added": 0})
            print(f"  [WARN] '{shape.name}': section(s) {missing} not in Excel.")
            return 0

    # ── Chart ────────────────────────────────────────────────────────────────
    if shape.has_chart:
        n, logs, items = process_chart(shape, spec, section_data)
        for l in logs:
            print(f"    {l}")
        if n > 0:
            report_rows.append({**base, "status": STATUS_OK,
                                "items_applied": "\n".join(items), "triangles_added": n})
        else:
            report_rows.append({**base, "status": STATUS_NO_SIG,
                                "items_applied": "", "triangles_added": 0})
        return n

    # ── Table ────────────────────────────────────────────────────────────────
    if shape.has_table:
        # Skip 1-row legend tables
        if len(shape.table.rows) <= 1:
            return 0
        n, logs, items = process_table(shape, spec, section_data)
        for l in logs:
            print(f"    {l}")
        if n > 0:
            report_rows.append({**base, "status": STATUS_OK,
                                "items_applied": "\n".join(items), "triangles_added": n})
        else:
            report_rows.append({**base, "status": STATUS_NO_SIG,
                                "items_applied": "", "triangles_added": 0})
        return n

    # ── Textbox ──────────────────────────────────────────────────────────────
    if shape.has_text_frame:
        code = spec["codes"][0]
        sec  = section_data.get(code, {})

        # Rank lookup mode
        if spec.get("rank") is not None:
            rank = spec["rank"]
            ranked = sec.get("ranked_w5", [])
            if rank < 1 or rank > len(ranked):
                report_rows.append({**base, "status": STATUS_WARN, "items_applied": "",
                                    "triangles_added": 0})
                print(f"  '{shape.name}' → rank {rank} out of range")
                return 0
            _, lbl, _, _, p4, p5, tri = ranked[rank - 1]
            if tri:
                write_triangle_to_textbox(shape, tri)
                print(f"  '{shape.name}' [{col_label}] rank{rank} '{lbl[:40]}' → {tri}")
                report_rows.append({**base, "status": STATUS_OK,
                                    "items_applied": f"rank {rank} '{lbl[:40]}' → {tri}{_pct_suffix((p4, p5))}",
                                    "triangles_added": 1})
                return 1
            print(f"  '{shape.name}' → rank {rank} ('{lbl[:30]}') no sig")
            report_rows.append({**base, "status": STATUS_NO_SIG, "items_applied": "",
                                "triangles_added": 0})
            return 0

        # Exact-label lookup
        label = spec.get("label")
        if label is None:
            report_rows.append({**base, "status": STATUS_WARN, "items_applied": "",
                                "triangles_added": 0})
            return 0
        tri = sec.get("by_label", {}).get(normalise(label))
        if tri:
            write_triangle_to_textbox(shape, tri)
            print(f"  '{shape.name}' [{col_label}] → {tri}")
            _pct = sec.get("pct_by_label", {}).get(normalise(label))
            report_rows.append({**base, "status": STATUS_OK,
                                "items_applied": f"'{label}' → {tri}{_pct_suffix(_pct)}",
                                "triangles_added": 1})
            return 1
        print(f"  '{shape.name}' → no sig")
        report_rows.append({**base, "status": STATUS_NO_SIG, "items_applied": "",
                            "triangles_added": 0})
        return 0

    # Unknown shape type
    report_rows.append({**base, "status": STATUS_WARN, "items_applied": "",
                        "triangles_added": 0})
    return 0


# ── Auto-scan PPT ─────────────────────────────────────────────────────────────

def scan_shapes(prs):
    """Return list of (slide_num, shape, cut_tuple)."""
    override_map = {}
    for entry in SHAPE_OVERRIDES:
        slide, name, cut = _override_cut(entry)
        override_map[(slide, name)] = cut
    default_cut = _default_cut()
    out = []
    for slide_num, slide in enumerate(prs.slides, 1):
        if slide_num in SKIP_SLIDES:
            continue
        for shape in slide.shapes:
            if not shape.name.startswith(SHAPE_PREFIX):
                continue
            cut = override_map.get((slide_num, shape.name), default_cut)
            out.append((slide_num, shape, cut))
    return out


# ── Progress report ───────────────────────────────────────────────────────────

def build_report(report_rows):
    """Build the progress report workbook and return it as xlsx bytes."""
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sync Report"
    ws.sheet_view.showGridLines = False

    H_FILL = PatternFill("solid", start_color="1F4E79")
    H_FONT = Font(color="FFFFFF", bold=True, name="Arial", size=10)
    FILLS  = {
        STATUS_OK:      PatternFill("solid", start_color="E2EFDA"),
        STATUS_NO_SIG:  PatternFill("solid", start_color="EDEDED"),
        STATUS_WARN:    PatternFill("solid", start_color="FCE4D6"),
    }
    C_FONT = Font(name="Arial", size=10)
    THIN   = Side(style="thin", color="BFBFBF")
    BDR    = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)

    COLS = [
        ("Slide", 8), ("Shape Name", 36), ("Status", 42), ("Items Applied", 70),
    ]
    for c, (h, w) in enumerate(COLS, 1):
        cell = ws.cell(row=1, column=c, value=h)
        cell.fill = H_FILL
        cell.font = H_FONT
        cell.alignment = Alignment(horizontal="center", vertical="center")
        cell.border = BDR
        ws.column_dimensions[get_column_letter(c)].width = w
    ws.row_dimensions[1].height = 20

    ws.merge_cells("A2:D2")
    banner = ws["A2"]
    ws.row_dimensions[2].height = 18

    for i, row in enumerate(report_rows, 3):
        status = row["status"]
        fill = FILLS.get(status, FILLS[STATUS_WARN])
        vals = [
            row["slide"], row["shape_name"], status, row.get("items_applied", ""),
        ]
        for c, val in enumerate(vals, 1):
            cell = ws.cell(row=i, column=c, value=val)
            cell.fill = fill
            cell.font = C_FONT
            cell.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
            cell.border = BDR
        ws.row_dimensions[i].height = 16

    ok  = sum(1 for r in report_rows if r["status"] == STATUS_OK)
    tri = sum(r.get("triangles_added", 0) for r in report_rows)
    default_label = " > ".join(_default_cut()) or "(none)"
    banner.value = (
        f"Sync complete — {len(report_rows)} shapes | {ok} with triangles | "
        f"{tri} triangles applied | Default: {default_label} | "
        f"Sheet: {EXCEL_SHEET}"
    )
    banner.font  = Font(name="Arial", size=10, bold=True, color="1F4E79")
    banner.fill  = PatternFill("solid", start_color="D9E1F2")
    banner.alignment = Alignment(horizontal="left", vertical="center")
    ws.freeze_panes = "A3"

    import io
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ── Main ──────────────────────────────────────────────────────────────────────

def run_sync(excel_bytes, pptx_bytes, log=print):
    """
    Core entry point. Takes Excel + PPTX as raw bytes, applies triangles,
    and returns a dict:
      {
        "pptx_bytes":   updated presentation as bytes,
        "report_bytes": progress report xlsx as bytes,
        "total":        number of triangles applied,
        "report_rows":  list of per-shape result dicts,
      }
    `log` is a callable used for progress messages (defaults to print).
    """
    import io

    # ── Load Excel ────────────────────────────────────────────────────────────
    wb = openpyxl.load_workbook(io.BytesIO(excel_bytes), data_only=True)
    if EXCEL_SHEET not in wb.sheetnames:
        raise ValueError(
            f"Sheet '{EXCEL_SHEET}' not found in the data file. "
            f"Available sheets: {wb.sheetnames}"
        )
    ws = wb[EXCEL_SHEET]
    log(f"Loaded Excel — sheet '{EXCEL_SHEET}'.")

    # Reset per-run caches that depend on the specific workbook/deck
    get_col_ancestry.__defaults__[0].clear()  # the _cache dict
    normalise.cache_clear()

    data_by_col, _all_sections = load_excel_data(ws, log=log)

    # ── Load PPT ──────────────────────────────────────────────────────────────
    prs = Presentation(io.BytesIO(pptx_bytes))
    load_theme_colors(prs)

    shapes = scan_shapes(prs)
    slide_lookup = {i + 1: s for i, s in enumerate(prs.slides)}
    log(f"Scanning {len(shapes)} x-shapes across {len(prs.slides)} slides "
        f"(skipping {SKIP_SLIDES if SKIP_SLIDES else 'none'}).")

    report_rows   = []
    total         = 0
    current_slide = None

    for slide_num, shape, cut in shapes:
        if slide_num != current_slide:
            log(f"--- Slide {slide_num} ---")
            current_slide = slide_num
        spec = parse_shape_name(shape.name)
        section_data = data_by_col[cut]
        n = process_shape(slide_num, slide_lookup[slide_num], shape, spec,
                          section_data, cut, report_rows)
        total += n

    # ── Serialise outputs to bytes ──────────────────────────────────────────────
    ppt_buf = io.BytesIO()
    prs.save(ppt_buf)
    ppt_buf.seek(0)

    report_bytes = build_report(report_rows)

    log(f"Done. {total} total triangle(s) applied.")
    return {
        "pptx_bytes":   ppt_buf.getvalue(),
        "report_bytes": report_bytes,
        "total":        total,
        "report_rows":  report_rows,
    }


def main():
    """CLI entry point — reads/writes files using the configured filenames."""
    xl_path  = SCRIPT_DIR / EXCEL_FILE
    ppt_path = SCRIPT_DIR / PPTX_FILE
    out_path = SCRIPT_DIR / OUTPUT_FILE
    rpt_path = SCRIPT_DIR / REPORT_FILE
    for p in (xl_path, ppt_path):
        if not p.exists():
            sys.exit(f"ERROR: File not found: {p}")

    result = run_sync(xl_path.read_bytes(), ppt_path.read_bytes(), log=print)
    out_path.write_bytes(result["pptx_bytes"])
    try:
        rpt_path.write_bytes(result["report_bytes"])
    except PermissionError:
        ts = datetime.now().strftime("%H%M%S")
        rpt_path = rpt_path.with_stem(rpt_path.stem + f"_{ts}")
        rpt_path.write_bytes(result["report_bytes"])
    print(f"Saved -> {out_path}")
    print(f"Report -> {rpt_path}")
    input("\nPress Enter to close...")


if __name__ == "__main__":
    main()